from pptx import Presentation
from pptx.util import Inches, Pt
import os
from PIL import Image

def create_presentation():
    # Initialize presentation with configurable layout parameters
    prs = Presentation()

    # ============= LAYOUT CONFIGURATION =============
    # Slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins and spacing
    SLIDE_MARGIN = Inches(0.5)
    TOP_IMAGE_SPACING = Inches(0.5)
    VERTICAL_SPACING = Inches(0.5)

    # Size ratios (relative to available content space)
    TOP_IMAGES_HEIGHT_RATIO = 0.5
    BOTTOM_IMAGE_HEIGHT_RATIO = 0.5
    BOTTOM_IMAGE_WIDTH_RATIO = 0.9

    # Text annotation settings
    TEXT_LEFT = SLIDE_MARGIN
    TEXT_TOP = Inches(0.1)
    TEXT_FONT_SIZE = Pt(16)
    TEXT_BOLD = True
    # ================================================

    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide_layout = prs.slide_layouts[6]

    png_files = sorted([f for f in os.listdir('.') if f.lower().endswith('.png')])

    content_width = SLIDE_WIDTH - 2*SLIDE_MARGIN
    content_height = SLIDE_HEIGHT - 2*SLIDE_MARGIN

    for i in range(0, len(png_files), 3):
        group = png_files[i:i+3]
        slide = prs.slides.add_slide(slide_layout)

        # Add slide title from first image's filename
        if group:
            first_image = group[0]
            title_text = extract_slide_title(first_image)
            add_slide_title(slide, title_text, TEXT_LEFT, TEXT_TOP, TEXT_FONT_SIZE, TEXT_BOLD)

        # Top images layout
        if len(group) >= 2:
            top_img_width = (content_width - TOP_IMAGE_SPACING) / 2
            top_img_height = content_height * TOP_IMAGES_HEIGHT_RATIO

            add_image(slide, group[0], SLIDE_MARGIN, SLIDE_MARGIN, top_img_width, top_img_height)
            add_image(slide, group[1], SLIDE_MARGIN + top_img_width + TOP_IMAGE_SPACING, 
                     SLIDE_MARGIN, top_img_width, top_img_height)

        # Bottom image layout
        if len(group) >= 3:
            bottom_img_width = content_width * BOTTOM_IMAGE_WIDTH_RATIO
            bottom_img_height = content_height * BOTTOM_IMAGE_HEIGHT_RATIO
            bottom_left = SLIDE_MARGIN + (content_width - bottom_img_width)/2
            bottom_top = SLIDE_HEIGHT - SLIDE_MARGIN - bottom_img_height

            add_image(slide, group[2], bottom_left, bottom_top, bottom_img_width, bottom_img_height)

    prs.save('ImagePresentation.pptx')

def extract_slide_title(filename):
    """Extract first 5 characters from filename stem"""
    stem = os.path.splitext(filename)[0]
    return stem[:6]

def add_slide_title(slide, text, left, top, font_size, bold):
    """Add text annotation to slide"""
    textbox = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.4))
    text_frame = textbox.text_frame
    text_frame.clear()  # Remove default empty paragraph

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.bold = bold
    p.font.size = font_size
    p.font.name = 'Calibri'  # Ensures consistent font across systems

def add_image(slide, img_path, left, top, max_width, max_height):
    """Add image with proper aspect ratio"""
    with Image.open(img_path) as img:
        img_ratio = img.width / img.height

    width = min(max_width, max_height * img_ratio)
    height = width / img_ratio

    if height > max_height:
        height = max_height
        width = height * img_ratio

    slide.shapes.add_picture(img_path, left, top, width, height)

if __name__ == "__main__":
    import os
    from fileutils import (
        data_dir,
        cd,
    )
    data_dir()
    cd(os.path.join('analysis', 'output', 'fixations_over_letters'))
    create_presentation()